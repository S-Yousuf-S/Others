"""
📚 Sensei – AI Study Companion
--------------------------------

Author: Yousuf S. R. Sakkaf
Status: In Development
Version: 4.0

Overview:
---------
Sensei is an AI-powered educational assistant designed to transform
static study materials into an interactive and explainable
learning experience.

Instead of passively reading notes, users can upload educational
content and interact with it conversationally through a tutor-like
AI interface powered by Google Gemini.

Sensei emphasizes clarity, transparency, and simplicity by combining
lightweight contextual retrieval with structured AI-generated responses.

🔍 Core Capabilities:
- Multi-format document ingestion:
    • PDF
    • Jupyter Notebook (.ipynb)
    • PowerPoint (.pptx)

- AI-powered contextual question answering
- Dynamic conversational interface with session persistence
- Dual-layer response architecture:
    • Tutor-style generated explanation
    • Retrieved source evidence expanders

- Explainable retrieval workflow
- Dojo Mode premium UI theme
- API key based Gemini integration
- Chunk-based lightweight retrieval system

🧠 System Workflow:
1. User uploads study materials
2. Text is extracted and serialized into JSON storage
3. Relevant contextual chunks are retrieved using:
    • keyword scoring
    • stop-word filtering
    • chunk ranking

4. Gemini generates:
    • a concise tutor-style explanation
    • source-supported contextual responses

5. Retrieved evidence is displayed transparently through
   expandable source containers.

🎯 Design Philosophy:
Sensei prioritizes:
- simplicity over over-engineering
- explainability over black-box responses
- usability over unnecessary complexity
- lightweight architecture over heavy infrastructure

The project intentionally avoids vector databases and embedding
pipelines during Phase 1 in order to maintain:
- simplicity
- maintainability
- faster deployment
- educational clarity

🎨 UI/UX & Branding:
Sensei follows a Neo-Traditional AI design philosophy combining:
- calmness
- mentorship
- intelligence
- minimalism
- futuristic aesthetics

Tagline:
    "Inspired by Tradition. Powered by Intelligence."

UI enhancements include:
- premium dark interface
- Dojo Mode visual theme
- custom branding system
- responsive sidebar controls
- premium chat styling
- explainable response expanders
- contextual visual hierarchy

🚀 Architectural Highlights:
- Lightweight RAG-inspired workflow
- Contextual chunk retrieval
- Explainable AI response pipeline
- Session-persistent conversational memory
- Streamlit-powered interactive UI
- Google Gemini integration

🚀 Future Enhancements:
- Vector database integration
- Semantic search using embeddings
- Multi-session conversations
- Excel document ingestion
- Advanced personalization
- SDK migration to `google-genai`
- Enhanced retrieval intelligence

📌 Educational Value:
This project demonstrates practical application of:
- Large Language Models (LLMs)
- Retrieval-Augmented Generation concepts
- Explainable AI workflows
- Human-centered AI interaction
- Lightweight AI system design

while balancing:
- usability
- transparency
- maintainability
- performance

# NOTE:
# Current implementation uses the deprecated
# `google.generativeai` SDK for compatibility
# during development and review stages.
#
# Future versions may migrate to the newer
# `google-genai` SDK.
"""

# -------------------------------
# 📦 IMPORTS
# -------------------------------
import streamlit as st
import google.generativeai as genai
from PyPDF2 import PdfReader
from pptx import Presentation
import json
import os
from PIL import Image

# -------------------------------
# 🔑 CONFIG
# -------------------------------

#genai.configure(api_key="AIzaSyBCuMR_v2JOS5RH6Z1GinoRRipBUZ-CXP8")

logo = Image.open("Sensei_v1.png")
# -------------------------------
# 🔐 API KEY SYSTEM
# -------------------------------

if "api_configured" not in st.session_state:
    st.session_state.api_configured = False

if "model" not in st.session_state:
    st.session_state.model = None


# -------------------------------
# 🧠 SESSION INIT
# -------------------------------
if "initialized" not in st.session_state:
    st.session_state.initialized = True
    if os.path.exists("chat.json"):
        os.remove("chat.json")

# -------------------------------
# 📄 FILE READERS
# -------------------------------
def read_pdf(file):
    reader = PdfReader(file)
    return "".join([p.extract_text() or "" for p in reader.pages])

def read_ipynb(file):
    data = json.load(file)
    return "\n".join(
        ["".join(cell.get("source", "")) for cell in data["cells"]]
    )

def read_pptx(file):
    prs = Presentation(file)
    text = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text.append(shape.text)
    return "\n".join(text)

# -------------------------------
# 🔍 RETRIEVAL
# -------------------------------
def find_relevant_sources(question, documents):

    results = []

    stop_words = {
        "what", "is", "are", "the", "how", "and",
        "of", "to", "in", "for", "a", "an",
        "can", "you", "explain", "with", "which", "tell", 
        "me", "summarise", "write", "about", "from",
        "list", "give", "define", "does"
    }

    keywords = [
        word.lower()
        for word in question.split()
        if word.lower() not in stop_words
    ]

    for source, text in documents.items():

        # Split into larger chunks instead of lines
        chunks = text.split("\n\n")

        scored_chunks = []

        for chunk in chunks:

            chunk_lower = chunk.lower()

            score = 0

            for keyword in keywords:

                if keyword in chunk_lower:
                    score += 1

            # Strong filtering
            if score >= 2:

                scored_chunks.append((score, chunk.strip()))

        # Sort by relevance score
        scored_chunks.sort(
            key=lambda x: x[0],
            reverse=True
        )

        top_chunks = [
            chunk
            for score, chunk in scored_chunks[:3]
        ]

        if top_chunks:

            combined = "\n\n".join(top_chunks)

            results.append((source, combined))

    return results
# -------------------------------
# 🤖 GEMINI
# -------------------------------
def generate_summary(question, relevant_sources):
    context = ""
    for _, text in relevant_sources:
        context += text[:1000] + "\n\n"

    prompt = f"""
    You are Sensei, an expert tutor.

    Provide a clear, concise explanation.
    Use simple language and examples.

    Context:
    {context}

    Question:
    {question}
    """

    return st.session_state.model.generate_content(prompt).text


def ask_tutor(question, context):
    prompt = f"""
    You are Sensei, a professional AI tutor.

    Explain clearly using bullet points and examples.
    Answer ONLY from the context.

    Context:
    {context}

    Question:
    {question}
    """

    return st.session_state.model.generate_content(prompt).text


# -------------------------------
# 🎨 UI
# -------------------------------

st.set_page_config(
    page_title="Sensei AI",
    page_icon=logo,
    layout="wide"
)

st.markdown("""
<style>

/* -------------------------------
   GLOBAL
--------------------------------*/

html, body, [class*="css"] {
    font-family: 'Segoe UI', sans-serif;
}

/* -------------------------------
   MAIN CONTAINER
--------------------------------*/

.block-container {
    padding-top: 2rem;
    padding-bottom: 2rem;
}

/* -------------------------------
   HEADER
--------------------------------*/

.main-header {
    padding: 1.5rem;
    border-radius: 20px;
    background: linear-gradient(
        135deg,
        rgba(40,40,60,0.95),
        rgba(25,25,35,0.95)
    );
    margin-bottom: 2rem;
    border: 1px solid rgba(255,255,255,0.08);
}

.logo-section {
    display: flex;
    align-items: center;
    gap: 20px;
}

.logo-circle {
    width: 70px;
    height: 70px;
    border-radius: 50%;
    background: linear-gradient(135deg, #4F46E5, #06B6D4);
    display: flex;
    justify-content: center;
    align-items: center;
    font-size: 32px;
    box-shadow: 0 0 20px rgba(79,70,229,0.4);
}

.main-header h1 {
    margin: 0;
    font-size: 2.8rem;
    color: white;
}

.main-header p {
    margin-top: 5px;
    color: #b0b0b0;
    font-size: 1rem;
}

/* -------------------------------
   CHAT INPUT
--------------------------------*/

.stChatInputContainer {
    border-radius: 20px;
}

/* -------------------------------
   EXPANDERS
--------------------------------*/

.streamlit-expanderHeader {
    font-size: 1rem;
    font-weight: 600;
}

/* -------------------------------
   SIDEBAR
--------------------------------*/

section[data-testid="stSidebar"] {
    background-color: #161A22;
    border-right: 1px solid rgba(255,255,255,0.08);
}

/* -------------------------------
   FOOTER
--------------------------------*/

.footer {
    margin-top: 4rem;
    text-align: center;
    color: #888;
    font-size: 0.9rem;
    padding: 1rem;
}

</style>
""", unsafe_allow_html=True)

# -------------------------------
# 📂 SIDEBAR
# -------------------------------
with st.sidebar:
    st.image(logo, width=70)
    st.markdown("## Sensei")
    st.caption("Inspired by Tradition. Powered by Intelligence.")
    st.markdown("""
    A next-generation AI learning companion designed
    to combine traditional mentorship principles with
    modern contextual intelligence.

    ---
    """)
    dojo_mode = st.toggle("🥷 Enable Dojo Mode")
    st.markdown("---")
    
    # -------------------------------
    # 🔑 GEMINI API
    # -------------------------------

    st.markdown("### 🔑 Gemini API")

    user_api_key = st.text_input(
        "Enter Gemini API Key",
        type="password",
        placeholder="Paste your API key here..."
    )

    if st.button("✅ Connect API"):

        if user_api_key.strip() == "":
            st.error("Please enter a valid API key.")

        else:
            try:
                genai.configure(api_key=user_api_key)

                st.session_state.model = genai.GenerativeModel(
                    "gemini-3-flash-preview"
                )

                st.session_state.api_configured = True

                st.success("Gemini API connected successfully.")

            except Exception as e:
                st.error("Invalid API key or connection failed.")


    uploaded_files = st.file_uploader(
        "Upload study materials",
        accept_multiple_files=True,
        key=st.session_state.get("uploader_key", "uploader_1")
    )

    if st.button("🧹 Clear Chat"):
        st.session_state.chat = []
        if os.path.exists("chat.json"):
            os.remove("chat.json")

    if st.button("🗑️ Clear Documents"):

        if os.path.exists("documents.json"):
            os.remove("documents.json")

        st.session_state.chat = []

        if os.path.exists("chat.json"):
            os.remove("chat.json")

        documents = {}
        
        import uuid
        st.session_state["uploader_key"] = str(uuid.uuid4())
        
        st.success("Documents cleared successfully.")

        st.rerun()

# -------------------------------
# 🏛️ HEADER
# -------------------------------

col1, col2 = st.columns([1, 5])

with col1:
    st.image(logo, width=100)

with col2:
    st.html("""
        <div style="padding-top:10px;">
            <h1 style="
                margin-bottom:0;
                color:#F4F4F5;
                letter-spacing:0.5px;
                font-size:48px;
            ">
                Sensei
            </h1>

            <p style="
                font-size:18px;
                opacity:0.82;
                color:#D1D5DB;
                margin-top:6px;
            ">
                Inspired by Tradition. Powered by Intelligence.
            </p>
        </div>
    """)


# -------------------------------
# 🎨 DOJO CSS
# -------------------------------
if dojo_mode:
    st.markdown(
        """
        <style>

        /* =========================================
           MAIN APP BACKGROUND
        ========================================= */

        .stApp {
            background:
                radial-gradient(circle at top left,
                    rgba(61,217,180,0.18),
                    transparent 22%
                ),

                radial-gradient(circle at top right,
                    rgba(255,215,0,0.10),
                    transparent 18%
                ),

                radial-gradient(circle at bottom right,
                    rgba(61,217,180,0.10),
                    transparent 20%
                ),

                linear-gradient(
                    135deg,
                    #03110C 0%,
                    #081712 30%,
                    #101820 65%,
                    #05080D 100%
                );

            color: #F4F4F5;
        }

        /* =========================================
           SIDEBAR
        ========================================= */

        section[data-testid="stSidebar"] {
            background:
                linear-gradient(
                    180deg,
                    #161B22 0%,
                    #11161C 100%
                );

            border-right: 1px solid rgba(61,217,180,0.18);
        }

        /* =========================================
           HEADER CARD
        ========================================= */

        .main-card {
            background:
                linear-gradient(
                    135deg,
                    rgba(25,30,35,0.96),
                    rgba(35,40,45,0.90)
                );

            border: 1px solid rgba(61,217,180,0.28);

            border-radius: 24px;

            padding: 28px;
            
            position: relative;
            overflow: hidden;
            
            box-shadow:
                0 0 30px rgba(61,217,180,0.18),
                0 0 70px rgba(212,175,55,0.08);
        }
    
        /* =========================================
           HERO HEADER
        ========================================= */

        [data-testid="stHorizontalBlock"] {
            background:
                linear-gradient(
                    135deg,
                    rgba(20,28,32,0.92),
                    rgba(15,22,28,0.88)
                );

            border: 1px solid rgba(61,217,180,0.16);

            border-radius: 24px;

            padding: 18px;

            margin-bottom: 24px;

            box-shadow:
                0 0 25px rgba(61,217,180,0.10),
                0 0 50px rgba(255,215,0,0.06);
        }
    
        .main-card::before {
            content: "";
            position: absolute;
            top: 0;
            left: -100%;
            width: 200%;
            height: 1px;
            background: linear-gradient(
                90deg,
                transparent,
                rgba(255,215,0,0.95),
                transparent
            );
        
            animation: shimmer 4s linear infinite;       
        }

        @keyframes shimmer {
            0% { left: -100%; }
            100% { left: 100%; }
        }

        /* =========================================
           CHAT INPUT
        ========================================= */

        .stChatInput textarea {
            background-color: #1A1F24 !important;

            color: #F4F4F5 !important;

            border-radius: 16px !important;

            border: 1px solid rgba(61,217,180,0.22) !important;

            padding: 12px !important;
        }

        /* =========================================
           BUTTONS
        ========================================= */

        .stButton button {
            background:
                linear-gradient(
                    135deg,
                    #3DD9B4,
                    #2CBFA0
                ) !important;

            color: #081018 !important;

            border: none !important;

            border-radius: 14px !important;

            font-weight: 700 !important;

            transition: all 0.25s ease;

            box-shadow:
                0 0 12px rgba(61,217,180,0.18);
        }

        .stButton button:hover {
            transform: translateY(-2px);

            box-shadow:
                0 0 20px rgba(61,217,180,0.32),
                0 0 40px rgba(255,215,0,0.20);
        }

        /* =========================================
           EXPANDERS
        ========================================= */

        .streamlit-expanderHeader {
            background: linear-gradient(
                90deg,
                rgba(61,217,180,0.08),
                rgba(212,175,55,0.08)
            );

            border-radius: 12px;

            border: 1px solid rgba(255,215,0,0.38);

            box-shadow:
                0 0 18px rgba(255,215,0,0.16);
        }

        /* =========================================
           CHAT MESSAGES
        ========================================= */

        .stChatMessage {

            background:
                linear-gradient(
                    135deg,
                    rgba(20,28,32,0.92),
                    rgba(15,22,28,0.88)
                );

            border-radius: 20px;

            padding: 16px;

            border: 1px solid rgba(61,217,180,0.12);

            backdrop-filter: blur(10px);

            box-shadow:
                0 0 18px rgba(61,217,180,0.08),
                0 0 30px rgba(255,215,0,0.05);

            margin-bottom: 12px;
        }
        /* =========================================
            TYPOGRAPHY
        ========================================= */

        h1, h2, h3 {
            letter-spacing: 0.5px;
        }

        /* =========================================
           SCROLLBAR
        ========================================= */

        ::-webkit-scrollbar {
            width: 8px;
        }

        ::-webkit-scrollbar-thumb {
            background:
                linear-gradient(
                    180deg,
                    #3DD9B4,
                    #D4AF37
                );

            border-radius: 20px;
        }

        /* =========================================
           HORIZONTAL LINE
        ========================================= */

        hr {
            border: none;
            height: 1px;

            background: linear-gradient(
                90deg,
                transparent,
                rgba(212,175,55,0.65),
                transparent
            );

            margin-top: 18px;
            margin-bottom: 18px;
        }

        </style>
        """,
        unsafe_allow_html=True
    )
# -------------------------------
# 📄 PROCESS FILES
# -------------------------------
documents = {}

if uploaded_files:
    for file in uploaded_files:
        if file.name.endswith(".pdf"):
            documents["pdf"] = documents.get("pdf", "") + read_pdf(file)
        elif file.name.endswith(".ipynb"):
            documents["notebook"] = read_ipynb(file)
        elif file.name.endswith(".pptx"):
            documents["slides"] = read_pptx(file)

    with open("documents.json", "w") as f:
        json.dump(documents, f)

if os.path.exists("documents.json"):
    with open("documents.json", "r") as f:
        documents = json.load(f)

# -------------------------------
# 💬 CHAT STATE
# -------------------------------
if "chat" not in st.session_state:
    if os.path.exists("chat.json"):
        with open("chat.json", "r") as f:
            st.session_state.chat = json.load(f)
    else:
        st.session_state.chat = []

# -------------------------------
# 💬 DISPLAY CHAT
# -------------------------------
for chat in st.session_state.chat:
    with st.chat_message("user"):
        st.markdown(chat["question"])

    with st.chat_message("assistant"):

        st.markdown("📌 **Main Explanation**")
        st.markdown(chat["summary"])
        st.markdown("---")

        for src, ans in chat["answers"]:

            label = ""
            if src == "pdf":
                label = "🟥 PDF Explanation"
            elif src == "notebook":
                label = "🟩 Notebook Code"
            elif src == "slides":
                label = "🟦 Slides Summary"

            with st.expander(label):
                st.caption(
                    "Relevant retrieved content used to generate the response."
                )
                if src == "notebook":
                    st.code(ans, language="python")
                else:
                    st.markdown(ans)

# -------------------------------
# ❓ INPUT
# -------------------------------
question = st.chat_input("Ask another question...")
if not st.session_state.api_configured:
    st.warning("Please connect your Gemini API key from the sidebar.")


if question and documents and st.session_state.api_configured:

    with st.chat_message("user"):
        st.markdown(question)

    relevant_sources = find_relevant_sources(question, documents)

    answers = []

    with st.chat_message("assistant"):
        with st.spinner("Sensei is thinking..."):

            if not relevant_sources:
                summary = "❌ Not found in uploaded material"
                st.markdown(summary)
            else:
                summary = generate_summary(question, relevant_sources)

                st.markdown("📌 **Main Explanation**")
                st.markdown(summary)
                st.markdown("---")

                for src, text in relevant_sources:
                    response = text
                    answers.append((src, response))

                    label = ""
                    if src == "pdf":
                        label = "🟥 PDF Explanation"
                    elif src == "notebook":
                        label = "🟩 Notebook Code"
                    elif src == "slides":
                        label = "🟦 Slides Summary"

                    with st.expander(label):
                        st.caption(
                            "Relevant retrieved content used to generate the response."
                        )
                        st.markdown(response)

    st.session_state.chat.append({
        "question": question,
        "summary": summary,
        "answers": answers
    })

    with open("chat.json", "w") as f:
        json.dump(st.session_state.chat, f)
    st.markdown("""
    <div class="footer">
        Built with Google Gemini • Designed by S.S.R. Yousuf
    </div>
    """, unsafe_allow_html=True)    
    