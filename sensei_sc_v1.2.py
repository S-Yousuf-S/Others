"""
📚 Sensei – Study Companion
--------------------------------

Author: Yousuf S. R. Sakkaf
Status: In-Progress
Version: 1.0

Overview:
---------
Sensei is an AI-powered educational assistant designed to transform
static study materials into an interactive learning experience.

Instead of passively reading notes, users can ask questions and receive
clear, structured, tutor-like explanations grounded strictly in their
uploaded content.

🔍 Core Capabilities:
- Multi-format document ingestion (PDF, Jupyter Notebook, PowerPoint)
- Context-aware Q&A using Google Gemini API
- Unified conversational interface with session persistence
- Dual-layer responses:
    • Concise summary (main explanation)
    • Source-wise breakdown (for transparency)

🧠 System Workflow:
1. User uploads study materials
2. Text is extracted and stored in JSON format
3. Relevant content is retrieved using keyword-based matching
4. Gemini generates:
    • A combined explanation (like a tutor)
    • Supporting answers per source

🎯 Design Philosophy:
- Simplicity over over-engineering (no vector DB in Phase 1)
- Explainability over black-box responses
- Clean and focused user experience

🚀 Future Enhancements:
- Multi-session chat system
- Vector database integration for scalable retrieval
- Semantic search using embeddings
- Advanced UI/UX improvements

This project demonstrates practical application of LLMs in education,
balancing usability, performance, and clarity.
"""

# -------------------------------
# 📦 IMPORTS
# -------------------------------
import streamlit as st
import google.generativeai as genai
from PyPDF2 import PdfReader
from pptx import Presentation
import json
import os
import time

# -------------------------------
# 🔑 CONFIG
# -------------------------------
genai.configure(api_key="AIzaSyBCuMR_v2JOS5RH6Z1GinoRRipBUZ-CXP8")
model = genai.GenerativeModel("gemini-2.5-flash")

# -------------------------------
# 🧠 SESSION INIT
# -------------------------------
if "initialized" not in st.session_state:
    st.session_state.initialized = True
    if os.path.exists("chat.json"):
        os.remove("chat.json")

# -------------------------------
# 📄 FILE READERS
# -------------------------------
def read_pdf(file):
    reader = PdfReader(file)
    return "".join([p.extract_text() or "" for p in reader.pages])

def read_ipynb(file):
    data = json.load(file)
    return "\n".join(
        ["".join(cell.get("source", "")) for cell in data["cells"]]
    )

def read_pptx(file):
    prs = Presentation(file)
    text = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text.append(shape.text)
    return "\n".join(text)

# -------------------------------
# 🔍 RETRIEVAL
# -------------------------------
def find_relevant_sources(question, documents):
    results = []
    keywords = question.lower().split()

    for source, text in documents.items():
        match_count = sum(word in text.lower() for word in keywords)

        if match_count >= 2:
            results.append((source, text[:1500]))

    return results

# -------------------------------
# 🤖 GEMINI
# -------------------------------
def generate_summary(question, relevant_sources):
    context = ""
    for _, text in relevant_sources:
        context += text[:1000] + "\n\n"

    prompt = f"""
    You are Sensei, an expert tutor.

    Provide a clear, concise explanation.
    Use simple language and examples.

    Context:
    {context}

    Question:
    {question}
    """

    return model.generate_content(prompt).text


def ask_tutor(question, context):
    prompt = f"""
    You are Sensei, a professional AI tutor.

    Return responses ONLY in proper markdown format.

    STRICT formatting rules:
    - Every bullet point MUST start on a new line
    - Use "-" for bullet points
    - Leave one empty line between paragraphs
    - Use markdown tables properly
    - Use headings where appropriate
    - Never combine multiple bullet points in one line
    - Keep formatting highly readable
    
    Answer ONLY from the context.

    Context:
    {context}

    Question:
    {question}
    """

    return model.generate_content(prompt).text


# -------------------------------
# 🎨 UI
# -------------------------------
st.set_page_config(page_title="Sensei AI", page_icon="📚", layout="wide")
st.markdown("""
<style>

.block-container {
    padding-top: 2rem;
}

</style>
""", unsafe_allow_html=True)

st.title("🧠 Sensei")
st.caption("Inspired by Tradition. Powered by Intelligence.")

# -------------------------------
# 📂 SIDEBAR
# -------------------------------
with st.sidebar:
    st.markdown("""
    ## 🧠 Sensei AI

    Inspired by Tradition. Powered by Intelligence.

    ---
    """)
    st.header("⚙️ Controls")

    uploaded_files = st.file_uploader(
        "Upload study materials",
        accept_multiple_files=True
    )

    if st.button("🧹 Clear Chat"):
        st.session_state.chat = []
        if os.path.exists("chat.json"):
            os.remove("chat.json")

    if st.button("🗑️ Clear Documents"):
        if os.path.exists("documents.json"):
            os.remove("documents.json")
        st.session_state.chat = []
        if os.path.exists("chat.json"):
            os.remove("chat.json")

# -------------------------------
# 📄 PROCESS FILES
# -------------------------------
documents = {}

if uploaded_files:
    for file in uploaded_files:
        if file.name.endswith(".pdf"):
            documents["pdf"] = documents.get("pdf", "") + read_pdf(file)
        elif file.name.endswith(".ipynb"):
            documents["notebook"] = read_ipynb(file)
        elif file.name.endswith(".pptx"):
            documents["slides"] = read_pptx(file)

    with open("documents.json", "w") as f:
        json.dump(documents, f)

if os.path.exists("documents.json"):
    with open("documents.json", "r") as f:
        documents = json.load(f)

# -------------------------------
# 💬 CHAT STATE
# -------------------------------
if "chat" not in st.session_state:
    if os.path.exists("chat.json"):
        with open("chat.json", "r") as f:
            st.session_state.chat = json.load(f)
    else:
        st.session_state.chat = []

# -------------------------------
# 💬 DISPLAY CHAT
# -------------------------------
for chat in st.session_state.chat:
    with st.chat_message("user"):
        st.markdown(chat["question"])

    with st.chat_message("assistant"):

        st.markdown("📌 **Main Explanation**")
        st.markdown(chat["summary"])
        st.markdown("---")

        for src, ans in chat["answers"]:

            label = ""
            if src == "pdf":
                label = "🟥 PDF Explanation"
            elif src == "notebook":
                label = "🟩 Notebook Code"
            elif src == "slides":
                label = "🟦 Slides Summary"

            with st.expander(label):
                if src == "notebook":
                    st.code(ans, language="python")
                else:
                    st.markdown(ans)

# -------------------------------
# ❓ INPUT
# -------------------------------
question = st.chat_input("Ask another question...")

if question and documents:

    with st.chat_message("user"):
        st.markdown(question)

    relevant_sources = find_relevant_sources(question, documents)

    answers = []

    with st.chat_message("assistant"):
        with st.spinner("Sensei is thinking..."):

            if not relevant_sources:
                summary = "❌ Not found in uploaded material"
                st.markdown(summary)
            else:
                summary = generate_summary(question, relevant_sources)
                summary = summary.replace("•", "\n- ")
                summary = summary.replace("* ", "\n* ")
                # Streaming effect
                st.markdown("📌 **Main Explanation**")
                st.markdown(summary)

                st.markdown("---")

                for src, text in relevant_sources:
                    response = text[:800]
                    formatted_response = response.replace("•", "\n- ")
                    formatted_response = formatted_response.replace("* ", "\n* ")
                    answers.append((src, formatted_response))

                    label = ""
                    if src == "pdf":
                        label = "🟥 PDF Explanation"
                    elif src == "notebook":
                        label = "🟩 Notebook Code"
                    elif src == "slides":
                        label = "🟦 Slides Summary"

                    with st.expander(label):
                        #if src == "notebook":
                        #    st.code(response, language="python")
                        #else:
                        st.markdown(response)
                        #formatted_response = response.replace("•", "\n- ")
                        #formatted_response = formatted_response.replace("* ", "\n* ")
                        #st.markdown(formatted_response)

    st.session_state.chat.append({
        "question": question,
        "summary": summary,
        "answers": answers
    })

    with open("chat.json", "w") as f:
        json.dump(st.session_state.chat, f)
st.markdown("---")
st.caption("Built with Google Gemini • Designed by S.S.R. Yousuf")        