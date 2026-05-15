"""
📚 Sensei – AI Study Companion
--------------------------------
Author          : Yousuf S. R. Sakkaf
Status          : Completed
Release Date    : 15/05/2026
Version         : 8.1

Overview:
---------
Sensei is an AI-powered educational assistant engineered to transform
static study materials into an interactive, conversational learning
experience. Powered by Google Gemini, it acts as a virtual tutor, 
allowing users to query their own notes and presentations with full
transparency.

Sensei employs a dual-layer response architecture: it generates concise,
tutor-style explanations while simultaneously providing the exact 
retrieved source evidence in dedicated UI expanders.

🔍 Core Capabilities:
- Multi-Format Ingestion: Natively parses PDF, Jupyter Notebook (.ipynb), and PowerPoint (.pptx).
- Transparent Retrieval: Displays exact file excerpts and code blocks used to formulate answers.
- Session Persistence: Maintains conversational context and document state per user session.
- Dojo Mode Interface: A premium, neo-traditional dark UI theme optimized for focused learning.

🧠 Architectural Philosophy (Phase 1):
This build intentionally bypasses heavy Vector Databases (like Chroma/Pinecone) 
and Embedding Pipelines in favor of a lightweight, JSON-serialized keyword 
retrieval system. This ensures:
    • Faster deployment and zero infrastructure overhead.
    • Maximum explainability (no black-box semantic mapping).
    • Easy code maintainability and high educational value regarding RAG fundamentals.

🚀 Future Roadmap:
- Migration to modern `google-genai` SDK.
- Vector database integration for semantic search.
- Expansion to Excel/Word document parsing.
"""

# -------------------------------
# 📦 IMPORTS
# -------------------------------
import streamlit as st
import google.generativeai as genai
from PyPDF2 import PdfReader
from pptx import Presentation
import json
import os
from PIL import Image
import re
import base64
from io import BytesIO

# -------------------------------
# 🔑 CONFIG
# -------------------------------

logo = Image.open(".\Assets\Sensei_v1.png")
# -------------------------------
# 🔐 API KEY SYSTEM
# -------------------------------

if "api_configured" not in st.session_state:
    st.session_state.api_configured = False

if "api_status" not in st.session_state:
    st.session_state.api_status = None

if "model" not in st.session_state:
    st.session_state.model = None


# -------------------------------
# 🧠 SESSION INIT
# -------------------------------
if "initialized" not in st.session_state:
    st.session_state.initialized = True
    # Generate a unique session ID to namespace files per user session
    import uuid as _uuid
    st.session_state.session_id = _uuid.uuid4().hex

# Derive per-session file paths
_sid = st.session_state.session_id
CHAT_FILE = f"chat_{_sid}.json"
DOCS_FILE = f"documents_{_sid}.json"

if "initialized" not in st.session_state or not os.path.exists(CHAT_FILE):
    pass  # fresh session — no stale chat to remove

# -------------------------------
# 📄 FILE READERS
# -------------------------------
def read_pdf(file):
    reader = PdfReader(file)
    return "".join([p.extract_text() or "" for p in reader.pages])

#def read_ipynb(file):
#    data = json.load(file)
#    return "\n".join(
#        ["".join(cell.get("source", "")) for cell in data["cells"]]
#    )

def read_ipynb(file):
    data = json.load(file)
    formatted_content = []
    
    for cell in data["cells"]:
        cell_type = cell.get("cell_type", "")
        source = "".join(cell.get("source", ""))
        
        if cell_type == "code":
            # Wrap code cells in markdown code blocks for beautiful rendering
            formatted_content.append(f"```python\n{source}\n```")
        else:
            # Leave markdown cells as they are
            formatted_content.append(source)
            
    return "\n\n".join(formatted_content)

def read_pptx(file):
    prs = Presentation(file)
    text = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text.append(shape.text)
    return "\n".join(text)

# -------------------------------
# 🧹 MARKDOWN NORMALIZER
# -------------------------------
def normalize_markdown(text):

    # Normalize line endings
    text = text.replace("\r\n", "\n")

    # Remove excessive blank lines (3+ → 2)
    #text = re.sub(r"\n{3,}", "\n\n", text)

    # Fix markdown heading spacing — only at line start, not inside code
    #text = re.sub(r'(?m)^(#{1,6})([^\s#])', r'\1 \2', text)

    # Strip trailing whitespace per line
    #text = "\n".join(line.rstrip() for line in text.splitlines())
    
    # NEW: Remove single newlines that are NOT followed by another newline
    # This stitches broken PDF sentences back together
    text = re.sub(r'(?<!\n)\n(?!\n)', ' ', text)

    # Remove excessive blank lines
    text = re.sub(r"\n{3,}", "\n\n", text)

    # Fix markdown heading spacing
    text = re.sub(r'(?m)^#([^\s])', r'# \1', text)
    return text.strip()


# -------------------------------
# 🧹 CHUNK FORMATTER FOR EXPANDERS
# -------------------------------
def format_chunks_for_display(raw_text, source_type):
    """
    Splits raw retrieved text into readable chunks and formats them
    with a visual divider. Deduplicates and cleans before rendering.
    """
    # Deduplicate lines first
    cleaned = clean_retrieved_content(raw_text)

    # Split into paragraph-level chunks
    chunks = [c.strip() for c in cleaned.split("\n\n") if c.strip()]

    # Filter out very short noise chunks (< 30 chars) unless it's a heading
    chunks = [
        c for c in chunks
        if len(c) >= 30 or c.startswith("#")
    ]

    return chunks

# -------------------------------
# 🖼️ RENDER NOTEBOOK MARKDOWN
# -------------------------------
def render_notebook_content(content):

    pattern = r'!\[.*?\]\((data:image\/.*?;base64,.*?)\)'

    matches = re.findall(pattern, content)

    # Remove images from markdown text
    #cleaned_text = re.sub(pattern, '', content)
    
    # Remove inline base64 images from markdown text so it doesn't clutter the screen
    cleaned_text = re.sub(pattern, '*[Image omitted from text]*', content)
    
    # SAFE: only add newline before markdown headings at line start,
    # not inside code blocks or Python comments
    #cleaned_text = re.sub(r'(?m)^(#{1,6} )', r'\n\1', cleaned_text)
    
    cleaned_text = re.sub(r'\n{3,}', '\n\n', cleaned_text)
    #cleaned_text = normalize_markdown(cleaned_text)

    # Render markdown FIRST
    st.markdown(cleaned_text)

    # Render extracted images separately
    for image_data in matches:

        try:
            header, encoded = image_data.split(",", 1)

            image_bytes = base64.b64decode(encoded)

            st.image(image_bytes)

        except:
            pass

# -------------------------------
# 🧼 RETRIEVED CONTENT CLEANER
# -------------------------------
def clean_retrieved_content(text):

    lines = text.splitlines()

    seen = set()

    cleaned = []

    for line in lines:

        normalized = line.strip().lower()

        if normalized == "":
            cleaned.append(line)
            continue

        if normalized not in seen:
            seen.add(normalized)
            cleaned.append(line)

    return "\n".join(cleaned)

# -------------------------------
# 🔍 RETRIEVAL
# -------------------------------
def find_relevant_sources(question, documents):

    results = []

    stop_words = {
        "what", "is", "are", "the", "how", "and",
        "of", "to", "in", "for", "a", "an",
        "can", "you", "explain", "with", "which", "tell", 
        "me", "summarise", "write", "about", "from",
        "list", "give", "define", "does"
    }

    #keywords = [
    #    word.lower()
    #    for word in question.split()
    #    if word.lower() not in stop_words
    #]

    # Extract only words, ignoring punctuation attached to them
    raw_words = re.findall(r'\b\w+\b', question.lower())
    
    keywords = [
        word for word in raw_words
        if word not in stop_words
    ]

    for source, text in documents.items():

        # Split into larger chunks instead of lines
        chunks = text.split("\n\n")

        scored_chunks = []

        for chunk_index, chunk in enumerate(chunks):

            chunk_lower = chunk.lower()

            # Truncate chunks longer than 1500 characters instead of skipping
            if len(chunk) > 1500:
                chunk = chunk[:1500]
                chunk_lower = chunk.lower()

            score = 0
            for keyword in keywords:

                if keyword in chunk_lower:
                    score += 1

            # Strong filtering
            if score >= 1:

                context_window = chunks[
                    max(0, chunk_index - 1):
                    min(len(chunks), chunk_index + 2)
                ]

                expanded_chunk = "\n\n".join(context_window)

                scored_chunks.append((score, expanded_chunk.strip()))

        # Sort by relevance score
        scored_chunks.sort(
            key=lambda x: x[0],
            reverse=True
        )

        top_chunks = [
            chunk
            for score, chunk in scored_chunks[:3]
        ]
        top_chunks = list(dict.fromkeys(top_chunks))
        if top_chunks:
            
            combined = "\n\n".join(top_chunks)

            results.append((source, combined))

    return results
# -------------------------------
# 🤖 GEMINI
# -------------------------------
def generate_summary(question, relevant_sources):
    context = ""
    for _, text in relevant_sources:
        context += text[:1000] + "\n\n"

    prompt = f"""
    You are Sensei, an expert tutor.

    Provide a clear, concise explanation.
    Use simple language and examples.

    Context:
    {context}

    Question:
    {question}
    """

    return st.session_state.model.generate_content(prompt).text



# -------------------------------
# 🎨 UI
# -------------------------------

st.set_page_config(
    page_title="Sensei AI",
    page_icon=logo,
    layout="wide"
)

st.markdown("""
<style>

/* -------------------------------
   GLOBAL
--------------------------------*/

html, body, [class*="css"] {
    font-family: 'Segoe UI', sans-serif;
}

/* -------------------------------
   MAIN CONTAINER
--------------------------------*/

.block-container {
    padding-top: 2rem;
    padding-bottom: 2rem;
}

/* -------------------------------
   HEADER
--------------------------------*/

.main-header {
    padding: 1.5rem;
    border-radius: 20px;
    background: linear-gradient(
        135deg,
        rgba(40,40,60,0.95),
        rgba(25,25,35,0.95)
    );
    margin-bottom: 2rem;
    border: 1px solid rgba(255,255,255,0.08);
}

.logo-section {
    display: flex;
    align-items: center;
    gap: 20px;
}

.logo-circle {
    width: 70px;
    height: 70px;
    border-radius: 50%;
    background: linear-gradient(135deg, #4F46E5, #06B6D4);
    display: flex;
    justify-content: center;
    align-items: center;
    font-size: 32px;
    box-shadow: 0 0 20px rgba(79,70,229,0.4);
}

.main-header h1 {
    margin: 0;
    font-size: 2.8rem;
    color: white;
}

.main-header p {
    margin-top: 5px;
    color: #b0b0b0;
    font-size: 1rem;
}

/* -------------------------------
   CHAT INPUT
--------------------------------*/

.stChatInputContainer {
    border-radius: 20px;
}

/* -------------------------------
   EXPANDERS
--------------------------------*/

.streamlit-expanderHeader {
    font-size: 1rem;
    font-weight: 600;
}

/* -------------------------------
   SIDEBAR
--------------------------------*/

section[data-testid="stSidebar"] {
    background-color: #161A22;
    border-right: 1px solid rgba(255,255,255,0.08);
}

/* -------------------------------
   FOOTER
--------------------------------*/

.footer {
    margin-top: 4rem;
    text-align: center;
    color: #888;
    font-size: 0.9rem;
    padding: 1rem;
}

</style>
""", unsafe_allow_html=True)

# -------------------------------
# 📂 SIDEBAR
# -------------------------------
with st.sidebar:
    st.image(logo, width=80)
    st.markdown("## Sensei")
    st.caption("Inspired by Tradition. Powered by Intelligence.")
    st.markdown("""
    <div class="upload-tip" style="color: #D1D5DB !important;">
    A next-generation AI learning companion designed
    to combine traditional mentorship principles with
    modern contextual intelligence.
    </div>
    """, unsafe_allow_html=True)

    dojo_mode = st.toggle("🥷 Enable Dojo Mode")
    st.markdown("---")
    
    # -------------------------------
    # 🔑 GEMINI API
    # -------------------------------

    st.markdown("### 🔑 Gemini API")

    user_api_key = st.text_input(
        "Enter Gemini API Key",
        type="password",
        placeholder="Paste your API key here..."
    )

    if st.button("✅ Connect API"):

        if user_api_key.strip() == "":
            st.session_state.api_status = "empty"

        else:
            try:
                genai.configure(api_key=user_api_key)

                test_model = genai.GenerativeModel(
                    "gemini-3-flash-preview"
                )

                # Real validation
                test_model.generate_content("Hello")

                st.session_state.model = test_model
                st.session_state.api_configured = True
                st.session_state.api_status = "success"

            except Exception:
                st.session_state.api_configured = False
                st.session_state.model = None
                st.session_state.api_status = "error"

    if st.session_state.api_status == "success":
        st.success("Gemini API connected successfully.")

    elif st.session_state.api_status == "error":
        st.error("Invalid API key or connection failed.")

    elif st.session_state.api_status == "empty":
        st.warning("Please enter a valid API key.")

    uploaded_files = st.file_uploader(
        "Upload study materials",
        accept_multiple_files=True,
        key=st.session_state.get("uploader_key", "uploader_1")
    )

    st.markdown("""
    <div class="upload-tip" style="color: #D1D5DB !important;">
        📚 For best retrieval quality, upload related study materials together.
    </div>
    """, unsafe_allow_html=True)

    if st.button("🧹 Clear Chat"):
        st.session_state.chat = []
        if os.path.exists(CHAT_FILE):
            os.remove(CHAT_FILE)

    if st.button("🗑️ Clear Documents"):

        if os.path.exists(DOCS_FILE):
            os.remove(DOCS_FILE)

        documents = {}

        import uuid
        st.session_state["uploader_key"] = str(uuid.uuid4())

        st.success("Documents cleared successfully.")

        st.rerun()

# -------------------------------
# 🏛️ HEADER
# -------------------------------

col1, col2 = st.columns([1, 5])

with col1:
    st.image(logo, width=120)

with col2:
    st.html("""
        <div style="padding-top:10px;">
            <h1 style="
                margin-bottom:0;
                color:#F4F4F5;
                letter-spacing:0.5px;
                font-size:48px;
            ">
                Sensei
            </h1>

            <p style="
                font-size:18px;
                opacity:0.82;
                color:#D1D5DB;
                margin-top:6px;
            ">
                Inspired by Tradition. Powered by Intelligence.
            </p>
        </div>
    """)


# -------------------------------
# 🎨 DOJO CSS
# -------------------------------
if dojo_mode:
    st.markdown(
        """
        <style>

        /* =========================================
           MAIN APP BACKGROUND
        ========================================= */

        .stApp {
            background:
                radial-gradient(circle at top left,
                    rgba(61,217,180,0.18),
                    transparent 22%
                ),

                radial-gradient(circle at top right,
                    rgba(255,215,0,0.10),
                    transparent 18%
                ),

                radial-gradient(circle at bottom right,
                    rgba(61,217,180,0.10),
                    transparent 20%
                ),

                linear-gradient(
                    135deg,
                    #03110C 0%,
                    #081712 30%,
                    #101820 65%,
                    #05080D 100%
                );

            color: #F4F4F5;
        }

        /* =========================================
           SIDEBAR
        ========================================= */

        section[data-testid="stSidebar"] {
            background:
                linear-gradient(
                    180deg,
                    #161B22 0%,
                    #11161C 100%
                );

            border-right: 1px solid rgba(61,217,180,0.18);
        }

        /* =========================================
           HEADER CARD
        ========================================= */

        .main-card {
            background:
                linear-gradient(
                    135deg,
                    rgba(25,30,35,0.96),
                    rgba(35,40,45,0.90)
                );

            border: 1px solid rgba(61,217,180,0.28);

            border-radius: 24px;

            padding: 28px;
            
            position: relative;
            overflow: hidden;
            
            box-shadow:
                0 0 30px rgba(61,217,180,0.18),
                0 0 70px rgba(212,175,55,0.08);
        }
    
        /* =========================================
           HERO HEADER
        ========================================= */

        [data-testid="stHorizontalBlock"] {
            background:
                linear-gradient(
                    135deg,
                    rgba(20,28,32,0.92),
                    rgba(15,22,28,0.88)
                );

            border: 1px solid rgba(61,217,180,0.16);

            border-radius: 24px;

            padding: 18px;

            margin-bottom: 24px;

            box-shadow:
                0 0 25px rgba(61,217,180,0.10),
                0 0 50px rgba(255,215,0,0.06);
        }
    
        .main-card::before {
            content: "";
            position: absolute;
            top: 0;
            left: -100%;
            width: 200%;
            height: 1px;
            background: linear-gradient(
                90deg,
                transparent,
                rgba(255,215,0,0.95),
                transparent
            );
        
            animation: shimmer 4s linear infinite;       
        }

        @keyframes shimmer {
            0% { left: -100%; }
            100% { left: 100%; }
        }

        /* =========================================
           CHAT INPUT
        ========================================= */

        .stChatInput textarea {
            background-color: #1A1F24 !important;

            color: #F4F4F5 !important;

            border-radius: 16px !important;

            border: 1px solid rgba(61,217,180,0.22) !important;

            padding: 12px !important;
        }

        /* =========================================
           BUTTONS
        ========================================= */

        .stButton button {
            background:
                linear-gradient(
                    135deg,
                    #3DD9B4,
                    #2CBFA0
                ) !important;

            color: #081018 !important;

            border: none !important;

            border-radius: 14px !important;

            font-weight: 700 !important;

            transition: all 0.25s ease;

            box-shadow:
                0 0 12px rgba(61,217,180,0.18);
        }

        .stButton button:hover {
            transform: translateY(-2px);

            box-shadow:
                0 0 20px rgba(61,217,180,0.32),
                0 0 40px rgba(255,215,0,0.20);
        }

        /* =========================================
           EXPANDERS
        ========================================= */

        .streamlit-expanderHeader {
            background: linear-gradient(
                90deg,
                rgba(61,217,180,0.08),
                rgba(212,175,55,0.08)
            );

            border-radius: 12px;

            border: 1px solid rgba(255,215,0,0.38);

            box-shadow:
                0 0 18px rgba(255,215,0,0.16);
        }

        /* =========================================
           CHAT MESSAGES
        ========================================= */

        .stChatMessage {

            background:
                linear-gradient(
                    135deg,
                    rgba(20,28,32,0.92),
                    rgba(15,22,28,0.88)
                );

            border-radius: 20px;

            padding: 16px;

            border: 1px solid rgba(61,217,180,0.12);

            backdrop-filter: blur(10px);

            box-shadow:
                0 0 18px rgba(61,217,180,0.08),
                0 0 30px rgba(255,215,0,0.05);

            margin-bottom: 12px;
        }
        /* =========================================
            TYPOGRAPHY
        ========================================= */

        h1, h2, h3 {
            letter-spacing: 0.5px;
        }

        /* =========================================
           SCROLLBAR
        ========================================= */

        ::-webkit-scrollbar {
            width: 8px;
        }

        ::-webkit-scrollbar-thumb {
            background:
                linear-gradient(
                    180deg,
                    #3DD9B4,
                    #D4AF37
                );

            border-radius: 20px;
        }

        /* =========================================
           HORIZONTAL LINE
        ========================================= */

        hr {
            border: none;
            height: 1px;

            background: linear-gradient(
                90deg,
                transparent,
                rgba(212,175,55,0.65),
                transparent
            );

            margin-top: 18px;
            margin-bottom: 18px;
        }

        /* =========================================
            SIDEBAR UPLOAD TIP
        ========================================= */

        .upload-tip, .upload-tip p, .upload-tip *{
            background:
                linear-gradient(
                    135deg,
                    rgba(61,217,180,0.10),
                    rgba(212,175,55,0.08)
                );

            border: 1px solid rgba(61,217,180,0.22);

            color: #D1D5DB !important;       

            box-shadow:
                0 0 12px rgba(61,217,180,0.08);

            font-style: normal;
        }

        </style>
        """,
        unsafe_allow_html=True
    )
# -------------------------------
# 📄 PROCESS FILES
# -------------------------------
documents = {}

if uploaded_files:
    for file in uploaded_files:
        if file.name.endswith(".pdf"):
            documents["pdf"] = documents.get("pdf", "") + read_pdf(file)
        elif file.name.endswith(".ipynb"):
            documents["notebook"] = documents.get("notebook", "") + "\n\n" + read_ipynb(file)
        elif file.name.endswith(".pptx"):
            documents["slides"] = documents.get("slides", "") + "\n\n" + read_pptx(file)

    with open(DOCS_FILE, "w") as f:
        json.dump(documents, f)

if os.path.exists(DOCS_FILE):
    with open(DOCS_FILE, "r") as f:
        documents = json.load(f)

# -------------------------------
# 💬 CHAT STATE
# -------------------------------
if "chat" not in st.session_state:
    if os.path.exists(CHAT_FILE):
        with open(CHAT_FILE, "r") as f:
            st.session_state.chat = json.load(f)
    else:
        st.session_state.chat = []

# -------------------------------
# 💬 DISPLAY CHAT
# -------------------------------
for chat in st.session_state.chat:
    with st.chat_message("user"):
        st.markdown(chat["question"])

    with st.chat_message("assistant"):

        st.markdown("📌 **Main Explanation**")
        st.markdown(chat["summary"])
        st.markdown("---")

        for src, ans in chat["answers"]:

            label = ""
            if src == "pdf":
                label = "🟥 PDF Source — Retrieved Passages"
            elif src == "notebook":
                label = "🟩 Notebook Source — Code & Notes"
            elif src == "slides":
                label = "🟦 Slides Source — Key Points"

            with st.expander(label, expanded=False):
                st.caption("📎 Relevant excerpts retrieved from your uploaded material.")
                st.divider()

                if src == "notebook":
                    render_notebook_content(ans)
                else:
                    chunks = format_chunks_for_display(ans, src)
                    for i, chunk in enumerate(chunks):
                        st.markdown(chunk)
                        if i < len(chunks) - 1:
                            st.markdown("---")

# -------------------------------
# ❓ INPUT
# -------------------------------
question = st.chat_input("Ask another question...")
if not st.session_state.api_configured:
    st.warning("Please connect your Gemini API key from the sidebar.")


if question and documents and st.session_state.api_configured:

    with st.chat_message("user"):
        st.markdown(question)

    relevant_sources = find_relevant_sources(question, documents)

    answers = []

    with st.chat_message("assistant"):

        if not relevant_sources:
            summary = "❌ Not found in uploaded material"
            st.markdown(summary)
        else:
            with st.spinner("Sensei is thinking..."):
                summary = generate_summary(question, relevant_sources)

            st.markdown("📌 **Main Explanation**")
            st.markdown(summary)
            st.markdown("---")

            for src, text in relevant_sources:
                response = text
                answers.append((src, response))

                label = ""
                if src == "pdf":
                    label = "🟥 PDF Source — Retrieved Passages"
                elif src == "notebook":
                    label = "🟩 Notebook Source — Code & Notes"
                elif src == "slides":
                    label = "🟦 Slides Source — Key Points"

                with st.expander(label, expanded=False):
                    st.caption("📎 Relevant excerpts retrieved from your uploaded material.")
                    st.divider()

                    if src == "notebook":
                        render_notebook_content(response)
                    else:
                        chunks = format_chunks_for_display(response, src)
                        for i, chunk in enumerate(chunks):
                            st.markdown(chunk)
                            if i < len(chunks) - 1:
                                st.markdown("---")

    st.session_state.chat.append({
        "question": question,
        "summary": summary,
        "answers": answers
    })

    with open(CHAT_FILE, "w") as f:
        json.dump(st.session_state.chat, f)

# -------------------------------
# 🔻 FOOTER
# -------------------------------
st.markdown("""
<div class="footer">
    Built with Google Gemini • Designed by S.S.R. Yousuf
</div>
""", unsafe_allow_html=True)