"""
📚 Sensei – Study Companion
--------------------------------

Author: Yousuf S. R. Sakkaf
Status: In-Progress
Version: 2.0

Overview:
---------
Sensei is an AI-powered educational assistant designed to transform
static study materials into an interactive learning experience.

Instead of passively reading notes, users can ask questions and receive
clear, structured, tutor-like explanations grounded strictly in their
uploaded content.

🔍 Core Capabilities:
- Multi-format document ingestion (PDF, Jupyter Notebook, PowerPoint)
- Context-aware Q&A using Google Gemini API
- Unified conversational interface with session persistence
- Dual-layer responses:
    • Concise summary (main explanation)
    • Source-wise breakdown (for transparency)

🧠 System Workflow:
1. User uploads study materials
2. Text is extracted and stored in JSON format
3. Relevant content is retrieved using keyword-based matching
4. Gemini generates:
    • A combined explanation (like a tutor)
    • Supporting answers per source

🎯 Design Philosophy:
- Simplicity over over-engineering (no vector DB in Phase 1)
- Explainability over black-box responses
- Clean and focused user experience

🚀 Future Enhancements:
- Multi-session chat system
- Vector database integration for scalable retrieval
- Semantic search using embeddings
- Advanced UI/UX improvements

This project demonstrates practical application of LLMs in education,
balancing usability, performance, and clarity.

------------------------
Version 2.0 Enhancements
------------------------

Sensei Version 2.0 introduces a refined user experience layer focused on
visual identity, usability, and modern educational interaction design.

🎨 UI/UX Enhancements:
- Introduced a premium dark-themed interface
- Added custom styling using Streamlit-compatible CSS
- Improved typography, spacing, and visual hierarchy
- Added collapsible response containers for cleaner navigation
- Enhanced notebook responses with syntax highlighting
- Added streaming-style tutor responses for conversational realism

🧠 Branding & Design Philosophy:
Sensei combines the discipline and guidance of traditional mentorship
with the intelligence and accessibility of modern AI systems.

Tagline:
    "Inspired by Tradition. Powered by Intelligence."

The visual identity follows a Neo-Traditional AI aesthetic:
- calm
- intelligent
- minimal
- mentor-inspired

🚀 Architectural Direction:
Version 2.0 intentionally focuses on front-end refinement while
preserving the lightweight backend architecture established in Phase 1.

This ensures:
- simplicity
- maintainability
- faster deployment
- improved user engagement

Future versions may include:
- vector database integration
- semantic search
- multi-session conversations
- advanced personalization

# NOTE:
# Current implementation uses the deprecated `google.generativeai`
# SDK for compatibility and simplicity during development.
# Future versions may migrate to the newer `google-genai` SDK.
"""

# -------------------------------
# 📦 IMPORTS
# -------------------------------
import streamlit as st
import google.generativeai as genai
from PyPDF2 import PdfReader
from pptx import Presentation
import json
import os
import time

# -------------------------------
# 🔑 CONFIG
# -------------------------------
genai.configure(api_key="AIzaSyBCuMR_v2JOS5RH6Z1GinoRRipBUZ-CXP8")
model = genai.GenerativeModel("gemini-2.5-flash")

# -------------------------------
# 🧠 SESSION INIT
# -------------------------------
if "initialized" not in st.session_state:
    st.session_state.initialized = True
    if os.path.exists("chat.json"):
        os.remove("chat.json")

# -------------------------------
# 📄 FILE READERS
# -------------------------------
def read_pdf(file):
    reader = PdfReader(file)
    return "".join([p.extract_text() or "" for p in reader.pages])

def read_ipynb(file):
    data = json.load(file)
    return "\n".join(
        ["".join(cell.get("source", "")) for cell in data["cells"]]
    )

def read_pptx(file):
    prs = Presentation(file)
    text = []
    for slide in prs.slides:
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text.append(shape.text)
    return "\n".join(text)

# -------------------------------
# 🔍 RETRIEVAL
# -------------------------------
def find_relevant_sources(question, documents):
    results = []
    keywords = question.lower().split()

    for source, text in documents.items():
        match_count = sum(word in text.lower() for word in keywords)

        if match_count >= 2:
            results.append((source, text[:1500]))

    return results

# -------------------------------
# 🤖 GEMINI
# -------------------------------
def generate_summary(question, relevant_sources):
    context = ""
    for _, text in relevant_sources:
        context += text[:1000] + "\n\n"

    prompt = f"""
    You are Sensei, an expert tutor.

    Provide a clear, concise explanation.
    Use simple language and examples.

    Context:
    {context}

    Question:
    {question}
    """

    return model.generate_content(prompt).text


def ask_tutor(question, context):
    prompt = f"""
    You are Sensei, a professional AI tutor.

    Explain clearly using bullet points and examples.
    Answer ONLY from the context.

    Context:
    {context}

    Question:
    {question}
    """

    return model.generate_content(prompt).text


# -------------------------------
# 🎨 UI
# -------------------------------
st.set_page_config(page_title="Sensei AI", page_icon="📚", layout="wide")
st.markdown("""
<style>

/* -------------------------------
   GLOBAL
--------------------------------*/

html, body, [class*="css"] {
    font-family: 'Segoe UI', sans-serif;
}

/* -------------------------------
   MAIN CONTAINER
--------------------------------*/

.block-container {
    padding-top: 2rem;
    padding-bottom: 2rem;
}

/* -------------------------------
   HEADER
--------------------------------*/

.main-header {
    padding: 1.5rem;
    border-radius: 20px;
    background: linear-gradient(
        135deg,
        rgba(40,40,60,0.95),
        rgba(25,25,35,0.95)
    );
    margin-bottom: 2rem;
    border: 1px solid rgba(255,255,255,0.08);
}

.logo-section {
    display: flex;
    align-items: center;
    gap: 20px;
}

.logo-circle {
    width: 70px;
    height: 70px;
    border-radius: 50%;
    background: linear-gradient(135deg, #4F46E5, #06B6D4);
    display: flex;
    justify-content: center;
    align-items: center;
    font-size: 32px;
    box-shadow: 0 0 20px rgba(79,70,229,0.4);
}

.main-header h1 {
    margin: 0;
    font-size: 2.8rem;
    color: white;
}

.main-header p {
    margin-top: 5px;
    color: #b0b0b0;
    font-size: 1rem;
}

/* -------------------------------
   CHAT INPUT
--------------------------------*/

.stChatInputContainer {
    border-radius: 20px;
}

/* -------------------------------
   EXPANDERS
--------------------------------*/

.streamlit-expanderHeader {
    font-size: 1rem;
    font-weight: 600;
}

/* -------------------------------
   SIDEBAR
--------------------------------*/

section[data-testid="stSidebar"] {
    background-color: #161A22;
    border-right: 1px solid rgba(255,255,255,0.08);
}

/* -------------------------------
   FOOTER
--------------------------------*/

.footer {
    margin-top: 4rem;
    text-align: center;
    color: #888;
    font-size: 0.9rem;
    padding: 1rem;
}

</style>
""", unsafe_allow_html=True)

st.markdown("""
<div class="main-header">
    <div class="logo-section">
        <div class="logo-circle">🧠</div>
        <div>
            <h1>Sensei</h1>
            <p>Inspired by Tradition. Powered by Intelligence.</p>
        </div>
    </div>
</div>
""", unsafe_allow_html=True)

# -------------------------------
# 📂 SIDEBAR
# -------------------------------
with st.sidebar:
    st.header("⚙️ Controls")
    st.markdown("""
    ### 🧠 Sensei AI
    A next-generation AI learning companion designed
    to combine traditional mentorship principles with
    modern contextual intelligence.

    ---
    """)
    uploaded_files = st.file_uploader(
        "Upload study materials",
        accept_multiple_files=True
    )

    if st.button("🧹 Clear Chat"):
        st.session_state.chat = []
        if os.path.exists("chat.json"):
            os.remove("chat.json")

    if st.button("🗑️ Clear Documents"):
        if os.path.exists("documents.json"):
            os.remove("documents.json")
        st.session_state.chat = []
        if os.path.exists("chat.json"):
            os.remove("chat.json")

# -------------------------------
# 📄 PROCESS FILES
# -------------------------------
documents = {}

if uploaded_files:
    for file in uploaded_files:
        if file.name.endswith(".pdf"):
            documents["pdf"] = documents.get("pdf", "") + read_pdf(file)
        elif file.name.endswith(".ipynb"):
            documents["notebook"] = read_ipynb(file)
        elif file.name.endswith(".pptx"):
            documents["slides"] = read_pptx(file)

    with open("documents.json", "w") as f:
        json.dump(documents, f)

if os.path.exists("documents.json"):
    with open("documents.json", "r") as f:
        documents = json.load(f)

# -------------------------------
# 💬 CHAT STATE
# -------------------------------
if "chat" not in st.session_state:
    if os.path.exists("chat.json"):
        with open("chat.json", "r") as f:
            st.session_state.chat = json.load(f)
    else:
        st.session_state.chat = []

# -------------------------------
# 💬 DISPLAY CHAT
# -------------------------------
for chat in st.session_state.chat:
    with st.chat_message("user"):
        st.markdown(chat["question"])

    with st.chat_message("assistant"):

        st.markdown("📌 **Main Explanation**")
        st.markdown(chat["summary"])
        st.markdown("---")

        for src, ans in chat["answers"]:

            label = ""
            if src == "pdf":
                label = "🟥 PDF Explanation"
            elif src == "notebook":
                label = "🟩 Notebook Code"
            elif src == "slides":
                label = "🟦 Slides Summary"

            with st.expander(label):
                if src == "notebook":
                    st.code(ans, language="python")
                else:
                    st.markdown(ans)

# -------------------------------
# ❓ INPUT
# -------------------------------
question = st.chat_input("Ask another question...")

if question and documents:

    with st.chat_message("user"):
        st.markdown(question)

    relevant_sources = find_relevant_sources(question, documents)

    answers = []

    with st.chat_message("assistant"):
        with st.spinner("Sensei is thinking..."):

            if not relevant_sources:
                summary = "❌ Not found in uploaded material"
                st.markdown(summary)
            else:
                summary = generate_summary(question, relevant_sources)

                # Streaming effect
                st.markdown("📌 **Main Explanation**")
                placeholder = st.empty()
                streamed_text = ""

                for word in summary.split():
                    streamed_text += word + " "
                    placeholder.markdown(streamed_text)
                    time.sleep(0.02)

                st.markdown("---")

                for src, text in relevant_sources:
                    response = ask_tutor(question, text)
                    answers.append((src, response))

                    label = ""
                    if src == "pdf":
                        label = "🟥 PDF Explanation"
                    elif src == "notebook":
                        label = "🟩 Notebook Code"
                    elif src == "slides":
                        label = "🟦 Slides Summary"

                    with st.expander(label):
                        if src == "notebook":
                            st.code(response, language="python")
                        else:
                            st.markdown(response)

    st.session_state.chat.append({
        "question": question,
        "summary": summary,
        "answers": answers
    })

    with open("chat.json", "w") as f:
        json.dump(st.session_state.chat, f)
    st.markdown("""
    <div class="footer">
        Built with Google Gemini • Designed by S.S.R. Yousuf
    </div>
    """, unsafe_allow_html=True)    
    